from __future__ import annotations

from io import BytesIO

from pptx import Presentation


def extract_segments(content: bytes) -> list[dict]:
    presentation = Presentation(BytesIO(content))
    segments: list[dict] = []
    for slide_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue
                segments.append(
                    {
                        "id": f"slide{slide_idx}_shape{shape.shape_id}_para{para_idx}",
                        "order": len(segments),
                        "source_text": text,
                        "plain_text": text,
                        "style_name": f"slide-{slide_idx}",
                        "segment_type": "paragraph",
                        "slide_index": slide_idx,
                        "shape_id": shape.shape_id,
                        "paragraph_index": para_idx,
                    }
                )
    return segments
