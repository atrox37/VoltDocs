from __future__ import annotations

from io import BytesIO

from pptx import Presentation


def export_pptx(original_bytes: bytes, segments: list[dict]) -> bytes:
    presentation = Presentation(BytesIO(original_bytes))
    by_id = {item["id"]: item for item in segments}
    for slide_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                key = f"slide{slide_idx}_shape{shape.shape_id}_para{para_idx}"
                item = by_id.get(key)
                if item:
                    paragraph.text = item.get("translation") or item.get("draftTranslation") or item.get("draft_translation") or ""
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
